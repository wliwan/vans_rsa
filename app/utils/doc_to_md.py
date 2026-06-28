"""文件转 Markdown 工具

支持格式: .txt .md .pdf .docx .ppt .pptx .xlsx .xls .csv
"""

import logging
import os

logger = logging.getLogger(__name__)


def file_to_markdown(file_path: str, filename: str) -> tuple[str, int]:
    """将文件内容转换为 Markdown 文本。返回 (markdown_text, char_count)"""
    ext = os.path.splitext(filename)[1].lower()
    if not ext:
        ext = os.path.splitext(file_path)[1].lower()

    if ext in (".txt", ".md", ".markdown"):
        return _read_text(file_path)
    if ext == ".pdf":
        return _pdf_to_md(file_path)
    if ext == ".docx":
        return _docx_to_md(file_path)
    if ext in (".ppt", ".pptx"):
        return _pptx_to_md(file_path)
    if ext in (".xlsx", ".xls"):
        return _excel_to_md(file_path)
    if ext == ".csv":
        return _csv_to_md(file_path)

    raise ValueError(f"不支持的文件格式: {ext}")


def _read_text(path: str) -> tuple[str, int]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    return text, len(text)


def _pdf_to_md(path: str) -> tuple[str, int]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    parts = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            parts.append(f"## 第 {i + 1} 页\n\n{text.strip()}")
    md = "\n\n---\n\n".join(parts)
    return md, len(md)


def _docx_to_md(path: str) -> tuple[str, int]:
    from docx import Document

    try:
        return _docx_to_md_via_python_docx(path)
    except Exception as e:
        logger.warning(f"python-docx 解析失败({e})，回退到 zip+xml 解析: {path}")
        try:
            return _docx_to_md_via_zip(path)
        except Exception as e2:
            raise RuntimeError(f"docx 解析完全失败: {e2}")


def _docx_to_md_via_python_docx(path: str) -> tuple[str, int]:
    from docx import Document

    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith("Heading"):
            level = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
            parts.append(f"{'#' * min(level, 4)} {text}")
        else:
            parts.append(text)
        parts.append("")

    for ti, table in enumerate(doc.tables):
        parts.append(f"\n### 表格 {ti + 1}\n")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
        parts.append("")

    md = "\n".join(parts)
    return md, len(md)


def _docx_to_md_via_zip(path: str) -> tuple[str, int]:
    """当 python-docx 遇到损坏 docx 时，回退到 zip + XML 直接提取文本"""
    import zipfile
    from xml.etree import ElementTree as ET

    # docx 的命名空间
    NS = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "t": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    parts = []

    with zipfile.ZipFile(path, "r") as zf:
        # 获取所有 XML 文件名，跳过无效条目
        xml_files = [
            name for name in zf.namelist()
            if name.endswith(".xml") and not name.startswith("_")
        ]

        for name in xml_files:
            try:
                data = zf.read(name)
            except Exception:
                continue

            try:
                root = ET.fromstring(data)
            except ET.ParseError:
                continue

            # 提取所有 <w:t> 中的文本
            texts = []
            for wt in root.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                if wt.text:
                    texts.append(wt.text)

            if texts:
                parts.append(" ".join(texts))

    md = "\n\n".join(parts)
    return md, len(md)


def _pptx_to_md(path: str) -> tuple[str, int]:
    from pptx import Presentation

    try:
        return _pptx_to_md_via_python_pptx(path)
    except Exception as e:
        logger.warning(f"python-pptx 解析失败({e})，回退到 zip+xml 解析: {path}")
        try:
            return _pptx_to_md_via_zip(path)
        except Exception as e2:
            raise RuntimeError(f"pptx 解析完全失败: {e2}")


def _pptx_to_md_via_python_pptx(path: str) -> tuple[str, int]:
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for si, slide in enumerate(prs.slides):
        parts.append(f"## 幻灯片 {si + 1}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
                parts.append("")
        parts.append("\n---\n")
    md = "\n".join(parts)
    return md, len(md)


def _pptx_to_md_via_zip(path: str) -> tuple[str, int]:
    """当 python-pptx 遇到损坏 pptx 时，回退到 zip + XML 直接提取文本"""
    import zipfile
    from xml.etree import ElementTree as ET

    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    parts = []

    with zipfile.ZipFile(path, "r") as zf:
        xml_files = [n for n in zf.namelist() if n.endswith(".xml")]
        for name in xml_files:
            try:
                data = zf.read(name)
            except Exception:
                continue
            try:
                root = ET.fromstring(data)
            except ET.ParseError:
                continue
            texts = []
            for t_elem in root.iter(f"{{{NS_A}}}t"):
                if t_elem.text:
                    texts.append(t_elem.text)
            if texts:
                parts.append(" ".join(texts))

    md = "\n\n".join(parts)
    return md, len(md)


def _excel_to_md(path: str) -> tuple[str, int]:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"## 工作表: {name}\n")
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        rows = rows[:200]
        headers = [str(c) if c is not None else "" for c in rows[0]]
        parts.append("| " + " | ".join(headers) + " |")
        parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            parts.append("| " + " | ".join(cells) + " |")
        parts.append("")
    wb.close()
    md = "\n".join(parts)
    return md, len(md)


def _csv_to_md(path: str) -> tuple[str, int]:
    import csv

    parts = []
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)[:200]
    if not rows:
        return "", 0
    headers = rows[0]
    parts.append("| " + " | ".join(headers) + " |")
    parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
    for row in rows[1:]:
        parts.append("| " + " | ".join(row) + " |")
    md = "\n".join(parts)
    return md, len(md)
